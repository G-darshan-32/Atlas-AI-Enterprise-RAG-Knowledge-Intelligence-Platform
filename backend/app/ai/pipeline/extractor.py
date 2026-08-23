"""Text extraction from various file formats."""
import io
from typing import Optional


def extract_text(content: bytes, file_type: str, filename: str = "") -> tuple[str, dict]:
    """
    Extract raw text and metadata from document bytes.
    Returns (text, metadata_dict).
    """
    extractors = {
        "pdf": _extract_pdf,
        "docx": _extract_docx,
        "pptx": _extract_pptx,
        "xlsx": _extract_xlsx,
        "markdown": _extract_markdown,
        "txt": _extract_txt,
        "csv": _extract_csv,
        "html": _extract_html,
        "ipynb": _extract_ipynb,
    }

    extractor = extractors.get(file_type, _extract_txt)
    return extractor(content, filename)


def _extract_pdf(content: bytes, filename: str) -> tuple[str, dict]:
    import fitz  # PyMuPDF
    doc = fitz.open(stream=content, filetype="pdf")
    pages = []
    for page in doc:
        pages.append(page.get_text())

    text = "\n\n".join(pages)

    # OCR fallback for scanned PDFs
    if len(text.strip()) < 100 and len(pages) > 0:
        text = _ocr_pdf(content)

    metadata = {
        "page_count": len(doc),
        "title": doc.metadata.get("title", ""),
        "author": doc.metadata.get("author", ""),
        "creation_date": doc.metadata.get("creationDate", ""),
    }
    doc.close()
    return text, metadata


def _ocr_pdf(content: bytes) -> str:
    """OCR fallback using pytesseract for scanned PDFs."""
    try:
        import fitz
        import pytesseract
        from PIL import Image

        doc = fitz.open(stream=content, filetype="pdf")
        text_parts = []
        for page in doc:
            pix = page.get_pixmap(dpi=200)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            text_parts.append(pytesseract.image_to_string(img))
        doc.close()
        return "\n\n".join(text_parts)
    except Exception:
        return ""


def _extract_docx(content: bytes, filename: str) -> tuple[str, dict]:
    from docx import Document
    doc = Document(io.BytesIO(content))

    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # Extract table content
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)

    text = "\n\n".join(paragraphs)
    metadata = {
        "author": doc.core_properties.author or "",
        "created": str(doc.core_properties.created or ""),
        "word_count": len(text.split()),
    }
    return text, metadata


def _extract_pptx(content: bytes, filename: str) -> tuple[str, dict]:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text = [f"[Slide {i}]"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        slides.append("\n".join(slide_text))

    text = "\n\n".join(slides)
    metadata = {"slide_count": len(prs.slides)}
    return text, metadata


def _extract_xlsx(content: bytes, filename: str) -> tuple[str, dict]:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
    sheets = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) for cell in row if cell is not None)
            if row_text.strip():
                rows.append(row_text)
        if rows:
            sheets.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))

    text = "\n\n".join(sheets)
    metadata = {"sheet_count": len(wb.worksheets)}
    return text, metadata


def _extract_markdown(content: bytes, filename: str) -> tuple[str, dict]:
    text = content.decode("utf-8", errors="replace")
    metadata = {"word_count": len(text.split())}
    return text, metadata


def _extract_txt(content: bytes, filename: str) -> tuple[str, dict]:
    import chardet
    detected = chardet.detect(content)
    encoding = detected.get("encoding") or "utf-8"
    text = content.decode(encoding, errors="replace")
    return text, {"encoding": encoding}


def _extract_csv(content: bytes, filename: str) -> tuple[str, dict]:
    import csv
    import chardet
    detected = chardet.detect(content)
    encoding = detected.get("encoding") or "utf-8"
    text_content = content.decode(encoding, errors="replace")
    reader = csv.reader(io.StringIO(text_content))
    rows = [" | ".join(row) for row in reader if any(cell.strip() for cell in row)]
    return "\n".join(rows), {"row_count": len(rows)}


def _extract_html(content: bytes, filename: str) -> tuple[str, dict]:
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(content, "html.parser")
    for tag in soup(["script", "style", "nav", "footer"]):
        tag.decompose()
    text = soup.get_text(separator="\n", strip=True)
    return text, {"title": soup.title.string if soup.title else ""}


def _extract_ipynb(content: bytes, filename: str) -> tuple[str, dict]:
    import json
    nb = json.loads(content.decode("utf-8", errors="replace"))
    cells = []
    for cell in nb.get("cells", []):
        cell_type = cell.get("cell_type", "")
        source = "".join(cell.get("source", []))
        if cell_type == "markdown":
            cells.append(source)
        elif cell_type == "code":
            cells.append(f"```python\n{source}\n```")
    return "\n\n".join(cells), {"cell_count": len(nb.get("cells", []))}
